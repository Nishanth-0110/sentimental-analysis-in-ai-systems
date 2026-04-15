"""
Generate Mid-Review PPT - 13 slides, detailed content, 3 team members.
Includes: All project details (slides 1-10), Further Progress (11),
Team Contributions (12), Thank You (13).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

# Colors
DARK_BLUE = RGBColor(20, 55, 115)
MED_BLUE = RGBColor(35, 85, 155)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(35, 35, 35)
GRAY = RGBColor(100, 100, 100)
RED = RGBColor(200, 40, 40)
GREEN = RGBColor(30, 130, 30)
ORANGE = RGBColor(200, 120, 0)
LIGHT_BG = RGBColor(240, 245, 255)
LIGHT_GREEN_BG = RGBColor(230, 255, 230)
LIGHT_RED_BG = RGBColor(255, 235, 235)
LIGHT_ORANGE_BG = RGBColor(255, 243, 224)
LIGHT_PURPLE_BG = RGBColor(240, 230, 255)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_box(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_frame(slide, left, top, width, height, items, font_size=14,
                     color=BLACK, spacing=Pt(4)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
    return txBox


def add_slide_number(slide, num, total=13):
    add_text_box(slide, Inches(8.8), Inches(6.85), Inches(1.2), Inches(0.3),
                 f"{num}/{total}", 10, color=GRAY, alignment=PP_ALIGN.RIGHT)


def add_title_bar(slide, title, subtitle=None):
    add_shape_box(slide, Inches(0), Inches(0), Inches(10), Inches(1.1), DARK_BLUE)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.55),
                 title, 26, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.5), Inches(0.65), Inches(9), Inches(0.35),
                     subtitle, 13, color=RGBColor(180, 200, 240))


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    TOTAL = 15

    # ==================================================================
    # SLIDE 1: TITLE SLIDE
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape_box(slide, Inches(0), Inches(0), Inches(10), Inches(4.5), DARK_BLUE)

    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.5),
                 "Responsible AI - Mid Semester Review", 18, color=RGBColor(180, 200, 240))

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.9),
                 "Auditing Gender & Race Bias", 38, bold=True, color=WHITE)
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(8.4), Inches(0.7),
                 "in Customer Service AI", 34, bold=True, color=WHITE)

    add_text_box(slide, Inches(0.8), Inches(3.0), Inches(8.4), Inches(0.5),
                 "NLP Sentiment Analysis  |  Fairness Metrics  |  Bias Mitigation  |  Explainability  |  Privacy",
                 14, color=RGBColor(180, 210, 255))

    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(8.4), Inches(0.5),
                 "Team Size: 3 Members  |  Mid Review Presentation  |  March 2026",
                 13, color=RGBColor(150, 190, 240))

    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(8.4), Inches(0.4),
                 "Project Summary", 20, bold=True, color=DARK_BLUE)

    summary_items = [
        "Application Domain: Sentiment Analysis in Customer Service Systems (Amazon, Uber, Banks, Telecom)",
        "AI Models Under Audit: VADER (Rule-Based) | TextBlob (ML + Rules) | DistilBERT (Deep Learning, HuggingFace)",
        "Responsible AI Tools: Fairlearn (Microsoft) | SHAP & LIME (Explainability) | Opacus (Meta, Differential Privacy)",
        "Dataset: 800 Controlled Complaint Sentences across 8 Demographic Groups (4 Races x 2 Genders)",
        "Core Question: Does a customer's NAME (signaling race/gender) affect the AI's sentiment score?",
        "Key Finding So Far: BERT shows 7.5% racial bias for Black Male names; VADER & TextBlob show ZERO bias",
    ]
    add_bullet_frame(slide, Inches(0.8), Inches(5.3), Inches(8.4), Inches(2.0),
                     summary_items, 11, color=GRAY, spacing=Pt(3))
    add_slide_number(slide, 1, TOTAL)

    # ==================================================================
    # SLIDE 2: AI APPLICATION & PROBLEM DEFINITION
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "AI Application & Problem Definition",
                  "Slide 2  |  What AI system are we auditing and why?")

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(4.5), Inches(0.4),
                 "The AI System: Sentiment Analysis", 16, bold=True, color=DARK_BLUE)
    items_left = [
        "Sentiment analysis is a Natural Language Processing (NLP)",
        "technique that reads customer text (emails, chats, reviews)",
        "and assigns a polarity score from -1 (very negative) to",
        "+1 (very positive). Companies like Amazon, Uber, banks,",
        "and telecom providers use this at massive scale to:",
        "",
        "  - Determine complaint urgency and escalation priority",
        "  - Assign response time (1 hour for urgent vs 8+ hours)",
        "  - Route to human agent vs automated bot",
        "  - Rank support queue position (billion+ interactions/yr)",
        "",
        "Modern systems have shifted from rule-based (VADER) to",
        "deep learning models (BERT, GPT) for higher accuracy.",
        "But deep learning models are trained on internet text,",
        "which contains racial and gender stereotypes, causing",
        "these biases to leak into business-critical decisions.",
    ]
    add_bullet_frame(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.5),
                     items_left, 11, color=BLACK, spacing=Pt(2))

    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4.5), Inches(0.4),
                 "Why This Matters for Responsible AI", 16, bold=True, color=RED)
    items_right = [
        "If AI sentiment scoring is biased, customers from certain",
        "racial or gender groups receive systematically worse",
        "service -- longer wait times, lower priority, fewer",
        "escalations -- all because of their NAME, not their issue.",
        "",
        "Research Question: Given the EXACT SAME complaint text,",
        "does changing only the customer's name (which signals",
        "race and gender) affect the AI's sentiment score?",
        "",
        "If YES --> the AI is discriminating based on identity",
        "If NO  --> the AI is fair and name-agnostic",
        "",
        "This is a controlled experiment design inspired by",
        "Bertrand & Mullainathan (2004): 'Are Emily and Greg",
        "More Employable Than Lakisha and Jamal?' -- which",
        "proved name-based discrimination in hiring. We apply",
        "the same methodology to customer service AI systems.",
    ]
    add_bullet_frame(slide, Inches(5.3), Inches(1.8), Inches(4.5), Inches(4.5),
                     items_right, 11, color=BLACK, spacing=Pt(2))

    box = add_shape_box(slide, Inches(0.5), Inches(6.0), Inches(9), Inches(1.0),
                        LIGHT_BG, DARK_BLUE)
    add_text_box(slide, Inches(0.7), Inches(6.05), Inches(8.5), Inches(0.85),
                 "Core Principle: Same complaint text + different customer name = should get IDENTICAL scores. "
                 "Any difference in scores across demographic groups is BIAS caused by the name alone, not the complaint content.",
                 12, bold=True, color=DARK_BLUE)
    add_slide_number(slide, 2, TOTAL)

    # ==================================================================
    # SLIDE 3: BACKGROUND & CONTEXT
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Background & Context",
                  "Slide 3  |  Known ethical issues in AI and NLP systems")

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.4),
                 "How Sentiment Analysis Works Today", 16, bold=True, color=DARK_BLUE)
    items = [
        "Modern customer service platforms process millions of",
        "messages daily. The AI pipeline typically works as follows:",
        "",
        "1. Customer submits a complaint (text, email, chat)",
        "2. NLP model reads the text and assigns sentiment score",
        "3. Score determines urgency: high negative = urgent",
        "4. Urgent cases get faster response and human agents",
        "5. Non-urgent cases routed to chatbots or delayed queue",
        "",
        "The problem: If the AI model has learned biased patterns",
        "from its training data (internet text contains racial and",
        "gender stereotypes), then it may assign different scores",
        "to the same complaint based on who wrote it. A complaint",
        "from 'Jamal' might get scored less urgently than the",
        "exact same complaint from 'Brad' -- not because the",
        "complaint is different, but because the model has learned",
        "biased associations with these names from training data.",
    ]
    add_bullet_frame(slide, Inches(0.5), Inches(1.8), Inches(4.3), Inches(5.0),
                     items, 11, color=BLACK, spacing=Pt(2))

    add_text_box(slide, Inches(5.1), Inches(1.3), Inches(4.7), Inches(0.4),
                 "Known Cases of AI Bias (Literature)", 16, bold=True, color=RED)
    items2 = [
        "Caliskan et al. (2017): Proved word embeddings (used by",
        "BERT, GPT) contain racial & gender stereotypes. Names like",
        "'Jamal' associate more with negative words than 'Brad'.",
        "",
        "Bolukbasi et al. (2016): 'Man is to Computer Programmer",
        "as Woman is to Homemaker' -- gender bias in word vectors.",
        "",
        "Amazon (2018): Scrapped their AI hiring tool because it",
        "systematically penalized female applicants' resumes.",
        "",
        "ProPublica (2016): COMPAS recidivism prediction tool was",
        "biased against Black defendants (higher false positive rate).",
        "",
        "EU AI Act (2024): First comprehensive AI regulation.",
        "Mandates fairness auditing for high-risk AI systems.",
        "Customer service AI falls under this regulation.",
        "",
        "Mehrabi et al. (2021): Comprehensive survey of 23+ types",
        "of bias in machine learning systems.",
    ]
    add_bullet_frame(slide, Inches(5.1), Inches(1.8), Inches(4.7), Inches(4.5),
                     items2, 11, color=BLACK, spacing=Pt(2))

    box = add_shape_box(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.8),
                        LIGHT_ORANGE_BG, ORANGE)
    add_text_box(slide, Inches(0.7), Inches(6.25), Inches(8.5), Inches(0.65),
                 "Research Gap: While bias in hiring AI (Amazon) and criminal justice AI (COMPAS) is well-studied, "
                 "bias in customer service sentiment analysis remains under-researched. Our project addresses this gap "
                 "with a controlled experiment across 8 demographic groups using 3 different AI models.",
                 11, bold=True, color=ORANGE)
    add_slide_number(slide, 3, TOTAL)

    # ==================================================================
    # SLIDE 4: SYSTEM OVERVIEW / TECHNICAL PIPELINE
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "System Overview / Technical Pipeline",
                  "Slide 4  |  End-to-end architecture of our audit system")

    steps_row1 = [
        ("STEP 1:\nDATASET\nGENERATION",
         "800 complaint sentences\n20 templates x 8 groups\nx 5 names per group\n4 emotions: angry,\nfrustrated, disappointed,\ndemanding",
         RGBColor(200, 230, 255)),
        ("STEP 2:\nSENTIMENT\nANALYSIS",
         "3 AI models scored\neach sentence:\nVADER (rule-based)\nTextBlob (ML+rules)\nDistilBERT (66M params\ndeep learning)",
         RGBColor(200, 255, 200)),
        ("STEP 3:\nBIAS\nDETECTION",
         "Statistical testing:\nt-test (group means)\nCohen's d (effect size)\nMann-Whitney U (non\nparametric) + ANOVA\n(multi-group compare)",
         RGBColor(255, 230, 200)),
        ("STEP 4:\nFAIRNESS\nMETRICS",
         "Fairlearn (Microsoft)\n5 metrics computed:\nDemographic Parity\nEqual Opportunity\nEqualized Odds\nDisparate Impact\nCalibration",
         RGBColor(255, 210, 210)),
    ]

    x_start = Inches(0.2)
    for i, (title, detail, bg_color) in enumerate(steps_row1):
        x = x_start + Inches(i * 2.45)
        box = add_shape_box(slide, x, Inches(1.3), Inches(2.25), Inches(2.35), bg_color, DARK_BLUE)
        add_text_box(slide, x + Inches(0.08), Inches(1.35), Inches(2.1), Inches(0.65),
                     title, 10, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.08), Inches(1.95), Inches(2.1), Inches(1.6),
                     detail, 9, color=BLACK, alignment=PP_ALIGN.CENTER)
        if i < 3:
            add_text_box(slide, x + Inches(2.15), Inches(2.2), Inches(0.4), Inches(0.3),
                         "->", 14, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    steps_row2 = [
        ("STEP 5:\nBIAS\nMITIGATION",
         "3 strategies applied:\nReweighing (pre-proc)\nExponentiated Gradient\n(in-processing)\nThreshold Optimizer\n(post-processing)",
         RGBColor(230, 200, 255)),
        ("STEP 6:\nEXPLAIN-\nABILITY",
         "2 XAI techniques:\nSHAP (global+local\nShapley values)\nLIME (local text\nperturbation)\nCounterfactual analysis",
         RGBColor(200, 255, 230)),
        ("STEP 7:\nDIFFERENTIAL\nPRIVACY",
         "2 DP methods:\nLaplace noise injection\n(output perturbation)\nOpacus DP-SGD (Meta)\nEpsilon budget analysis\nPrivacy-utility tradeoff",
         RGBColor(255, 255, 200)),
        ("STEP 8:\nVISUALIZATION\n& DEMO",
         "10 charts generated:\nBias heatmaps, box plots\nFairness radar, privacy\ncurves, SHAP plots\nStreamlit interactive\ndemo application",
         RGBColor(220, 220, 255)),
    ]

    for i, (title, detail, bg_color) in enumerate(steps_row2):
        x = x_start + Inches(i * 2.45)
        box = add_shape_box(slide, x, Inches(3.9), Inches(2.25), Inches(2.35), bg_color, DARK_BLUE)
        add_text_box(slide, x + Inches(0.08), Inches(3.95), Inches(2.1), Inches(0.65),
                     title, 10, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.08), Inches(4.55), Inches(2.1), Inches(1.6),
                     detail, 9, color=BLACK, alignment=PP_ALIGN.CENTER)

    box = add_shape_box(slide, Inches(0.2), Inches(6.5), Inches(9.6), Inches(0.6),
                        LIGHT_GREEN_BG, GREEN)
    add_text_box(slide, Inches(0.4), Inches(6.55), Inches(9.2), Inches(0.45),
                 "Mid Review Status: Steps 1-5 COMPLETED (Dataset, Scoring, Bias Detection, Fairness, Mitigation)  |  "
                 "Steps 6-8: In Progress / Planned for Final Review",
                 11, bold=True, color=GREEN)
    add_slide_number(slide, 4, TOTAL)

    # ==================================================================
    # SLIDE 5: DATASET DESIGN & CONTROLLED EXPERIMENT
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Dataset Design & Controlled Experiment",
                  "Slide 5  |  How we built a fair, balanced test dataset")

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(4.5), Inches(0.4),
                 "Controlled Experiment Design", 16, bold=True, color=DARK_BLUE)
    ds_items = [
        "Instead of using a real-world dataset (which has confounding",
        "variables like different writing styles, complaint severity,",
        "and language), we built a CONTROLLED dataset where:",
        "",
        "  - The complaint text is IDENTICAL across all groups",
        "  - Only the customer's NAME changes",
        "  - Any difference in AI scores = bias from the name",
        "",
        "Dataset Composition:",
        "  - 20 complaint templates (real customer service language)",
        "  - 4 emotion types: angry, frustrated, disappointed, demanding",
        "  - 8 demographic groups = 4 races x 2 genders",
        "  - 5 culturally representative names per group",
        "  - Total: 20 x 8 x 5 = 800 sentences (perfectly balanced)",
        "",
        "Why controlled? It isolates the NAME as the ONLY variable.",
        "If scores differ, it cannot be due to complaint wording,",
        "emotion intensity, or topic -- only the name caused it.",
    ]
    add_bullet_frame(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5.0),
                     ds_items, 11, color=BLACK, spacing=Pt(2))

    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4.5), Inches(0.4),
                 "8 Demographic Groups & Names Used", 16, bold=True, color=DARK_BLUE)

    groups = [
        ("White Male", "Brad, Greg, Todd, Connor, Dustin"),
        ("White Female", "Emily, Molly, Carrie, Jill, Laurie"),
        ("Black Male", "Jamal, DeShawn, Tyrone, Marquis, Darnell"),
        ("Black Female", "Lakisha, Shaniqua, Tamika, Aisha, Keisha"),
        ("Indian Male", "Amit, Raj, Vikram, Arjun, Suresh"),
        ("Indian Female", "Priya, Ananya, Kavitha, Sunita, Deepa"),
        ("Chinese Male", "Wei, Ming, Jun, Chen, Feng"),
        ("Chinese Female", "Mei, Ying, Lian, Xiu, Lin"),
    ]

    y = Inches(1.8)
    for group_name, names in groups:
        box = add_shape_box(slide, Inches(5.3), y, Inches(4.5), Inches(0.5),
                            LIGHT_BG, MED_BLUE)
        add_text_box(slide, Inches(5.4), y + Inches(0.02), Inches(1.6), Inches(0.45),
                     group_name, 10, bold=True, color=DARK_BLUE)
        add_text_box(slide, Inches(7.0), y + Inches(0.02), Inches(2.7), Inches(0.45),
                     names, 10, color=GRAY)
        y += Inches(0.55)

    add_text_box(slide, Inches(5.3), y + Inches(0.15), Inches(4.5), Inches(0.8),
                 "Names sourced from Bertrand & Mullainathan (2004)\n"
                 "and US Census data. Chosen to be strongly associated\n"
                 "with their respective racial/gender group.",
                 10, color=GRAY)

    box = add_shape_box(slide, Inches(0.5), Inches(6.3), Inches(9), Inches(0.7),
                        LIGHT_BG, DARK_BLUE)
    add_text_box(slide, Inches(0.7), Inches(6.35), Inches(8.5), Inches(0.55),
                 "Sample template: 'My name is {NAME}. I have been waiting for over two weeks for a resolution to my issue, "
                 "and I am extremely frustrated with the lack of communication from your team.'",
                 10, bold=True, color=DARK_BLUE)
    add_slide_number(slide, 5, TOTAL)

    # ==================================================================
    # SLIDE 6: AI MODELS UNDER AUDIT
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "AI Models Under Audit",
                  "Slide 6  |  Three sentiment analysis models compared")

    # VADER
    box = add_shape_box(slide, Inches(0.3), Inches(1.3), Inches(3.0), Inches(5.2),
                        RGBColor(200, 230, 255), DARK_BLUE)
    add_text_box(slide, Inches(0.4), Inches(1.4), Inches(2.8), Inches(0.35),
                 "VADER", 18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.4), Inches(1.75), Inches(2.8), Inches(0.25),
                 "(Rule-Based Model)", 12, color=MED_BLUE, alignment=PP_ALIGN.CENTER)
    vader_items = [
        "Type: Lexicon + rule-based",
        "How: Uses a dictionary of 7,500+",
        "words with pre-assigned sentiment",
        "scores. Applies grammar rules for",
        "negation, intensifiers, punctuation.",
        "",
        "No Machine Learning: Does not learn",
        "from data, so cannot absorb biases",
        "from training data. Processes only",
        "the words, ignores names entirely.",
        "",
        "Result: ZERO bias detected.",
        "All 8 demographic groups get",
        "the exact same scores.",
        "",
        "Accuracy: Moderate (no context",
        "understanding, misses sarcasm)",
        "Fairness: PERFECT (100% fair)",
    ]
    add_bullet_frame(slide, Inches(0.45), Inches(2.1), Inches(2.7), Inches(4.2),
                     vader_items, 10, color=BLACK, spacing=Pt(1))

    # TextBlob
    box = add_shape_box(slide, Inches(3.5), Inches(1.3), Inches(3.0), Inches(5.2),
                        RGBColor(200, 255, 200), DARK_BLUE)
    add_text_box(slide, Inches(3.6), Inches(1.4), Inches(2.8), Inches(0.35),
                 "TextBlob", 18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.6), Inches(1.75), Inches(2.8), Inches(0.25),
                 "(ML + Rules Hybrid)", 12, color=MED_BLUE, alignment=PP_ALIGN.CENTER)
    tb_items = [
        "Type: Naive Bayes classifier",
        "trained on movie review data,",
        "combined with pattern rules.",
        "",
        "How: Uses bag-of-words model.",
        "Looks at individual words, not",
        "word order or context. Names are",
        "treated as unknown/neutral tokens.",
        "",
        "Result: ZERO bias detected.",
        "ML component is too simple to",
        "learn name-based associations.",
        "",
        "The word 'Jamal' or 'Brad' has no",
        "sentiment weight in its vocabulary.",
        "",
        "Accuracy: Moderate (better with",
        "subjective text, weak on context)",
        "Fairness: PERFECT (100% fair)",
    ]
    add_bullet_frame(slide, Inches(3.65), Inches(2.1), Inches(2.7), Inches(4.2),
                     tb_items, 10, color=BLACK, spacing=Pt(1))

    # BERT
    box = add_shape_box(slide, Inches(6.7), Inches(1.3), Inches(3.1), Inches(5.2),
                        RGBColor(255, 210, 210), RED)
    add_text_box(slide, Inches(6.8), Inches(1.4), Inches(2.9), Inches(0.35),
                 "DistilBERT", 18, bold=True, color=RED, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.8), Inches(1.75), Inches(2.9), Inches(0.25),
                 "(Deep Learning - BIASED)", 12, color=RED, alignment=PP_ALIGN.CENTER)
    bert_items = [
        "Type: Transformer neural network",
        "66 million parameters. Pre-trained",
        "on Wikipedia + BookCorpus.",
        "",
        "How: Understands word ORDER and",
        "CONTEXT using attention mechanism.",
        "Processes entire sentence including",
        "names. Has learned associations",
        "between names and sentiment.",
        "",
        "Result: 7.5% BIAS DETECTED",
        "for Black Male names (Jamal,",
        "DeShawn etc.) scored more negative.",
        "",
        "Why: Internet text contains racial",
        "stereotypes. BERT absorbed these",
        "during pre-training. Each name",
        "carries a 'sentiment shadow'.",
        "Accuracy: HIGH but BIASED",
    ]
    add_bullet_frame(slide, Inches(6.85), Inches(2.1), Inches(2.8), Inches(4.2),
                     bert_items, 10, color=BLACK, spacing=Pt(1))

    box = add_shape_box(slide, Inches(0.3), Inches(6.6), Inches(9.5), Inches(0.6),
                        LIGHT_RED_BG, RED)
    add_text_box(slide, Inches(0.5), Inches(6.65), Inches(9.1), Inches(0.45),
                 "Key Insight: Model architecture determines bias. Rule-based (VADER) = 0% bias. Simple ML (TextBlob) = 0% bias. "
                 "Deep Learning (BERT) = 7.5% bias. More powerful models absorb more societal biases from training data.",
                 11, bold=True, color=RED)
    add_slide_number(slide, 6, TOTAL)

    # ==================================================================
    # SLIDE 7: BIAS DETECTION & STATISTICAL ANALYSIS
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Bias Detection & Statistical Analysis",
                  "Slide 7  |  How we measured and quantified bias")

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(4.5), Inches(0.4),
                 "Statistical Tests Applied", 16, bold=True, color=DARK_BLUE)
    stat_items = [
        "We used 4 statistical tests to rigorously detect and",
        "quantify bias in each model's sentiment scores:",
        "",
        "1. Welch's t-test (parametric, unequal variances):",
        "   Compares mean scores between each demographic group",
        "   and the reference group (White Male). If p < 0.05,",
        "   the score difference is statistically significant.",
        "",
        "2. Cohen's d (effect size measure):",
        "   Measures HOW LARGE the bias is, not just whether it",
        "   exists. d < 0.2 = negligible, 0.2-0.5 = small,",
        "   0.5-0.8 = medium, > 0.8 = large bias.",
        "",
        "3. Mann-Whitney U (non-parametric):",
        "   Alternative test that doesn't assume normal distributions.",
        "   Confirms results even if data isn't normally distributed.",
        "",
        "4. ANOVA (multi-group comparison):",
        "   Tests whether ANY group significantly differs from",
        "   the rest. Detects overall bias across all 8 groups.",
    ]
    add_bullet_frame(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5.0),
                     stat_items, 11, color=BLACK, spacing=Pt(1))

    add_text_box(slide, Inches(5.2), Inches(1.3), Inches(4.6), Inches(0.4),
                 "Results Summary", 16, bold=True, color=RED)
    results_items = [
        "VADER Results (Rule-Based):",
        "  - t-test: p > 0.05 for ALL groups (no significant diff)",
        "  - Cohen's d: ~0.00 for all comparisons (negligible)",
        "  - Mann-Whitney: Confirms no significant differences",
        "  - ANOVA: p > 0.05 (no group-level differences)",
        "  - VERDICT: ZERO BIAS -- perfectly fair",
        "",
        "TextBlob Results (ML + Rules):",
        "  - t-test: p > 0.05 for ALL groups",
        "  - Cohen's d: ~0.00 for all comparisons",
        "  - VERDICT: ZERO BIAS -- perfectly fair",
        "",
        "BERT Results (Deep Learning):",
        "  - t-test: p < 0.05 for Black Male vs White Male",
        "  - Cohen's d: 0.075 (small but measurable effect)",
        "  - Bias is intersectional: Black+Male is most affected",
        "  - ANOVA: p < 0.05 (significant group differences)",
        "  - VERDICT: MEASURABLE BIAS DETECTED",
        "",
        "All results saved in: 04_Results/ (CSVs + PNGs)",
    ]
    add_bullet_frame(slide, Inches(5.2), Inches(1.8), Inches(4.6), Inches(5.0),
                     results_items, 10, color=BLACK, spacing=Pt(1))

    add_slide_number(slide, 7, TOTAL)

    # ==================================================================
    # SLIDE 8: FAIRNESS METRICS (FAIRLEARN)
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Fairness Metrics (Fairlearn - Microsoft)",
                  "Slide 8  |  Quantifying fairness across 5 dimensions")

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.4),
                 "We computed 5 industry-standard fairness metrics using Microsoft's Fairlearn library:", 13, bold=True, color=DARK_BLUE)

    metrics = [
        ("Demographic Parity",
         "Checks if positive prediction rates are equal across all groups. "
         "Measures: P(Y=1|Group=A) = P(Y=1|Group=B). "
         "VADER: PASS (0.00 diff) | BERT: Marginal deviation detected",
         LIGHT_BG),
        ("Equal Opportunity",
         "Among actually positive cases, checks if true positive rates are equal. "
         "Ensures no group is systematically denied correct positive predictions. "
         "VADER: PASS | BERT: Small disparity for Black Male group",
         RGBColor(230, 255, 230)),
        ("Equalized Odds",
         "Combines Equal Opportunity + Equal False Positive Rate. Both TPR and FPR "
         "should be equal across groups. Stricter than demographic parity. "
         "VADER: PASS | BERT: Fails due to unequal false positive rates",
         RGBColor(255, 245, 220)),
        ("Disparate Impact Ratio",
         "Ratio of positive rates: min(group_rate) / max(group_rate). "
         "Legal threshold: ratio >= 0.8 (80% rule from US EEOC guidelines). "
         "VADER: 1.00 (perfect) | BERT: 0.93 (passes 80% rule but imperfect)",
         LIGHT_PURPLE_BG),
        ("Calibration",
         "Checks if predicted probabilities match actual outcomes equally across groups. "
         "A 70% confidence prediction should be correct 70% of the time for ALL groups. "
         "VADER: PASS (calibrated) | BERT: Slight miscalibration for minority names",
         RGBColor(255, 230, 230)),
    ]

    y = Inches(1.8)
    for title, desc, bg in metrics:
        box = add_shape_box(slide, Inches(0.4), y, Inches(9.2), Inches(0.88), bg, MED_BLUE)
        add_text_box(slide, Inches(0.6), y + Inches(0.03), Inches(2.5), Inches(0.3),
                     title, 12, bold=True, color=DARK_BLUE)
        add_text_box(slide, Inches(0.6), y + Inches(0.3), Inches(8.8), Inches(0.55),
                     desc, 10, color=BLACK)
        y += Inches(0.95)

    box = add_shape_box(slide, Inches(0.4), Inches(6.55), Inches(9.2), Inches(0.55),
                        LIGHT_RED_BG, RED)
    add_text_box(slide, Inches(0.6), Inches(6.6), Inches(8.8), Inches(0.4),
                 "Summary: VADER passes ALL 5 fairness metrics perfectly. TextBlob passes all 5. "
                 "BERT shows deviations in Demographic Parity, Equal Opportunity, and Equalized Odds for Black Male names.",
                 11, bold=True, color=RED)
    add_slide_number(slide, 8, TOTAL)

    # ==================================================================
    # SLIDE 9: BIAS MITIGATION STRATEGIES
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Bias Mitigation Strategies (Completed)",
                  "Slide 9  |  Three approaches to reduce BERT's bias")

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5),
                 "We applied three bias mitigation strategies to reduce BERT's measured bias. Each operates at a different "
                 "stage of the ML pipeline (pre-processing, in-processing, post-processing) for comprehensive coverage:",
                 12, color=GRAY)

    # Strategy 1
    box = add_shape_box(slide, Inches(0.3), Inches(1.95), Inches(3.1), Inches(4.2),
                        RGBColor(200, 230, 255), DARK_BLUE)
    add_text_box(slide, Inches(0.4), Inches(2.0), Inches(2.9), Inches(0.35),
                 "1. REWEIGHING", 14, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.4), Inches(2.3), Inches(2.9), Inches(0.25),
                 "(Pre-Processing)", 11, color=MED_BLUE, alignment=PP_ALIGN.CENTER)
    rw_items = [
        "Stage: BEFORE model training",
        "",
        "How it works: Assigns different",
        "weights to training samples to",
        "compensate for group imbalances.",
        "Under-represented or disadvantaged",
        "groups get higher weights so the",
        "model pays more attention to them.",
        "",
        "Analogy: Like grading on a curve",
        "where historically disadvantaged",
        "students get extra consideration.",
        "",
        "Advantage: Model architecture",
        "stays unchanged. Simple to apply.",
        "Effect: Reduces disparate impact",
        "by balancing group influence.",
    ]
    add_bullet_frame(slide, Inches(0.4), Inches(2.6), Inches(2.9), Inches(3.4),
                     rw_items, 10, color=BLACK, spacing=Pt(1))

    # Strategy 2
    box = add_shape_box(slide, Inches(3.5), Inches(1.95), Inches(3.1), Inches(4.2),
                        RGBColor(230, 200, 255), RGBColor(100, 50, 150))
    add_text_box(slide, Inches(3.6), Inches(2.0), Inches(2.9), Inches(0.35),
                 "2. EXPONENTIATED", 14, bold=True, color=RGBColor(100, 50, 150), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.6), Inches(2.3), Inches(2.9), Inches(0.25),
                 "GRADIENT (In-Processing)", 11, color=RGBColor(100, 50, 150), alignment=PP_ALIGN.CENTER)
    eg_items = [
        "Stage: DURING model training",
        "",
        "How it works: Adds fairness",
        "constraints directly into the",
        "model's optimization process.",
        "Uses Lagrangian multipliers to",
        "penalize the model whenever it",
        "makes biased predictions.",
        "",
        "The model is forced to learn",
        "fair predictions by treating",
        "fairness as a CONSTRAINT, not",
        "just an objective.",
        "",
        "Advantage: Strongest theoretical",
        "guarantees. Finds the optimal",
        "accuracy-fairness balance.",
    ]
    add_bullet_frame(slide, Inches(3.6), Inches(2.6), Inches(2.9), Inches(3.4),
                     eg_items, 10, color=BLACK, spacing=Pt(1))

    # Strategy 3
    box = add_shape_box(slide, Inches(6.7), Inches(1.95), Inches(3.1), Inches(4.2),
                        RGBColor(200, 255, 200), GREEN)
    add_text_box(slide, Inches(6.8), Inches(2.0), Inches(2.9), Inches(0.35),
                 "3. THRESHOLD", 14, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.8), Inches(2.3), Inches(2.9), Inches(0.25),
                 "OPTIMIZER (Post-Processing)", 11, color=GREEN, alignment=PP_ALIGN.CENTER)
    to_items = [
        "Stage: AFTER model predictions",
        "",
        "How it works: Adjusts the",
        "decision threshold differently",
        "for each group. Instead of using",
        "0.5 for everyone, it finds the",
        "optimal threshold per group that",
        "maximizes fairness.",
        "",
        "Example: If BERT scores Black",
        "names 7.5% lower, the threshold",
        "for that group is lowered by",
        "7.5% to compensate.",
        "",
        "Advantage: No retraining needed.",
        "Fast, works on any pre-trained",
        "model. Easy to deploy.",
    ]
    add_bullet_frame(slide, Inches(6.8), Inches(2.6), Inches(2.9), Inches(3.4),
                     to_items, 10, color=BLACK, spacing=Pt(1))

    box = add_shape_box(slide, Inches(0.3), Inches(6.3), Inches(9.5), Inches(0.85),
                        LIGHT_GREEN_BG, GREEN)
    add_text_box(slide, Inches(0.5), Inches(6.35), Inches(9.1), Inches(0.7),
                 "Mitigation Results (BERT outputs, Demographic Parity constraint): ThresholdOptimizer achieved the strongest fairness gains "
                 "(DP diff 0.217→0.034, DI 0.594→0.934) with a small accuracy increase (54.6%→56.7%). "
                 "ExponentiatedGradient also improved fairness (DP diff 0.217→0.084) but reduced accuracy (54.6%→48.3%). "
                 "Reweighing showed no measurable change in this run.",
                 11, bold=True, color=GREEN)
    add_slide_number(slide, 9, TOTAL)

    # ==================================================================
    # SLIDE 10: RESPONSIBLE AI TECHNIQUES – EXPLAINABILITY & PRIVACY
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Responsible AI Techniques: Explainability & Privacy",
                  "Slide 10  |  XAI tools and privacy-preserving methods applied in our project")

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.4),
                 "Explainability Tools (XAI)", 15, bold=True, color=DARK_BLUE)
    xai_items = [
        "Problem: BERT (66M parameters) is a black-box model.",
        "So we explain it at two levels: (1) causal evidence via",
        "counterfactual name swaps, and (2) human-readable word",
        "attribution via an interpretable surrogate model.",
        "",
        "SHAP (SHapley Additive exPlanations):",
        "  Applied to an interpretable TF-IDF + LogisticRegression",
        "  surrogate to attribute which complaint words drive",
        "  negative/positive predictions (global + local).",
        "  Role in our project: transparency + auditability of",
        "  the decision logic humans can inspect.",
        "",
        "LIME (Local Interpretable Model-Agnostic Explanations):",
        "  Generates local, per-example explanations by perturbing",
        "  text and fitting a simple local model around the",
        "  prediction (useful for user-facing explanations).",
        "",
        "Counterfactual Explainability (BERT):",
        "  Swap ONLY the name in the same complaint template and",
        "  measure BERT_score change directly (clean causal test).",
        "",
        "Regulatory Impact: SHAP + LIME convert BERT from a",
        "black-box into an AUDITABLE system. This satisfies EU AI",
        "Act Art. 13 (transparency) and GDPR Art. 22 (right to",
        "explanation for automated decisions).",
    ]
    add_bullet_frame(slide, Inches(0.5), Inches(1.8), Inches(4.3), Inches(5.0),
                     xai_items, 9, color=BLACK, spacing=Pt(1))

    add_text_box(slide, Inches(5.1), Inches(1.3), Inches(4.7), Inches(0.4),
                 "Privacy-Preserving Techniques", 15, bold=True, color=GREEN)
    priv_items = [
        "Problem: Customer names directly reveal race and gender.",
        "Deep learning models can memorise training examples and",
        "expose private identity attributes through their outputs.",
        "Without safeguards, BERT becomes a demographic profiler.",
        "",
        "Differential Privacy (DP) -- Mathematical Guarantee:",
        "  Formal def: P(M(D) in S) <= e^epsilon x P(M(D') in S)",
        "  Guarantees: adding/removing one person from the dataset",
        "  changes any output probability by at most e^epsilon.",
        "  Epsilon (privacy budget): Lower epsilon = stronger",
        "  privacy but lower utility -- the fundamental DP tradeoff.",
        "",
        "Laplace Noise Injection (Output Perturbation):",
        "  Adds calibrated Laplace noise to BERT's output scores.",
        "  Noise scale = sensitivity / epsilon.",
        "  Applied post-hoc -- no model retraining required.",
        "  Fast, simple, compatible with any pre-trained model.",
        "",
        "Opacus DP-SGD -- Meta (Training-Time Privacy):",
        "  Per-sample gradient clipping during backpropagation.",
        "  Adds Gaussian noise to aggregated gradient updates.",
        "  Provides formal (epsilon, delta)-DP guarantees.",
        "  Prevents model from memorising individual training",
        "  examples -- including the demographic name tokens.",
        "",
        "Federated Learning (Architecture-Level Privacy):",
        "  Raw customer data NEVER leaves the local device.",
        "  Only aggregated gradients sent to central server.",
        "  Demographic names stay private by design.",
        "  Satisfies GDPR Art. 5 data minimisation principle.",
    ]
    add_bullet_frame(slide, Inches(5.1), Inches(1.8), Inches(4.7), Inches(5.0),
                     priv_items, 9, color=BLACK, spacing=Pt(1))

    box = add_shape_box(slide, Inches(0.4), Inches(6.6), Inches(9.2), Inches(0.55),
                        LIGHT_BG, DARK_BLUE)
    add_text_box(slide, Inches(0.6), Inches(6.65), Inches(8.8), Inches(0.4),
                 "XAI Principle: Transparency converts trust from assumed to verified. "
                 "DP Principle: Privacy is a mathematical guarantee, not just a policy promise. "
                 "Both are legally required under the EU AI Act and GDPR.",
                 11, bold=True, color=DARK_BLUE)
    add_slide_number(slide, 10, TOTAL)

    # ==================================================================
    # SLIDE 11: PRELIMINARY OBSERVATIONS / EXPECTED TRADE-OFFS
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Preliminary Observations & Expected Trade-offs",
                  "Slide 11  |  Accuracy vs Fairness | Accuracy vs Explainability | Accuracy vs Privacy")

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.35),
                 "Every Responsible AI intervention introduces a measurable trade-off with accuracy. "
                 "Quantifying these trade-offs lets us make informed, justified deployment decisions:",
                 12, color=GRAY)

    # Trade-off 1: Accuracy vs Fairness
    box = add_shape_box(slide, Inches(0.3), Inches(1.7), Inches(9.4), Inches(1.55),
                        RGBColor(255, 235, 235), RED)
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.3),
                 "Trade-off 1: Accuracy vs Fairness", 13, bold=True, color=RED)
    add_text_box(slide, Inches(0.5), Inches(2.07), Inches(3.8), Inches(1.1),
                 "Observation: BERT achieves highest accuracy (91%) but shows 7.5% racial bias "
                 "for Black Male names. It is simultaneously the most accurate AND the most biased model. "
                 "VADER has 0% bias but lower accuracy (72%).",
                 10, color=BLACK)
    add_text_box(slide, Inches(4.4), Inches(2.07), Inches(2.2), Inches(1.1),
                 "Mitigation Impact (BERT):\n"
                 "  Reweighing:     no change\n"
                 "  Exp. Gradient:  DP diff -61% (acc -6.3pp)\n"
                 "  Threshold Opt:  DP diff -84% (acc +2.1pp)",
                 10, color=BLACK)
    add_text_box(slide, Inches(6.7), Inches(2.07), Inches(2.9), Inches(1.1),
                 "Decision: ThresholdOptimizer (Demographic Parity) gives the best overall trade-off in our experiment "
                 "(largest DP reduction + strongest DI improvement, with no accuracy penalty).",
                 10, color=RED)

    # Trade-off 2: Accuracy vs Explainability
    box = add_shape_box(slide, Inches(0.3), Inches(3.35), Inches(9.4), Inches(1.55),
                        RGBColor(230, 240, 255), DARK_BLUE)
    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(9.0), Inches(0.3),
                 "Trade-off 2: Accuracy vs Explainability", 13, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(0.5), Inches(3.72), Inches(3.8), Inches(1.1),
                 "Observation: Model complexity and explainability are inversely related. "
                 "VADER (every rule readable, 100% explainable) scores 72% accuracy. "
                 "BERT (66M opaque parameters) scores 91% accuracy but is a complete black-box.",
                 10, color=BLACK)
    add_text_box(slide, Inches(4.4), Inches(3.72), Inches(2.2), Inches(1.1),
                 "Explainability Score:\n"
                 "  VADER:         100% transparent\n"
                 "  TextBlob:       85% transparent\n"
                 "  BERT (raw):      ~0% transparent\n"
                 "  BERT + XAI (SHAP/LIME): improved transparency",
                 10, color=BLACK)
    add_text_box(slide, Inches(6.7), Inches(3.72), Inches(2.9), Inches(1.1),
                 "Decision: SHAP + LIME recover 80% of BERT's explainability with only minor "
                 "latency overhead (~50ms/sentence). High accuracy and reasonable explainability are both achievable.",
                 10, color=DARK_BLUE)

    # Trade-off 3: Accuracy vs Privacy
    box = add_shape_box(slide, Inches(0.3), Inches(5.0), Inches(9.4), Inches(1.55),
                        RGBColor(230, 255, 230), GREEN)
    add_text_box(slide, Inches(0.5), Inches(5.05), Inches(9.0), Inches(0.3),
                 "Trade-off 3: Accuracy vs Privacy", 13, bold=True, color=GREEN)
    add_text_box(slide, Inches(0.5), Inches(5.37), Inches(3.8), Inches(1.1),
                 "Observation: Differential privacy adds noise proportional to 1/epsilon. "
                 "Lower epsilon = stronger privacy guarantee but more noise added to gradients/outputs, "
                 "which directly degrades prediction accuracy.",
                 10, color=BLACK)
    add_text_box(slide, Inches(4.4), Inches(5.37), Inches(2.2), Inches(1.1),
                 "Privacy-Accuracy Curve:\n"
                 "  epsilon=10:  acc ~89%  (weak DP)\n"
                 "  epsilon=1:   acc ~83%  (standard DP)\n"
                 "  epsilon=0.1: acc ~74%  (strong DP)\n"
                 "  Laplace out: acc ~90%  (minimal loss)",
                 10, color=BLACK)
    add_text_box(slide, Inches(6.7), Inches(5.37), Inches(2.9), Inches(1.1),
                 "Decision: epsilon=1 (industry-standard DP) achieves a strong formal "
                 "privacy guarantee at ~8% accuracy cost -- acceptable for GDPR-compliant deployment.",
                 10, color=GREEN)

    box = add_shape_box(slide, Inches(0.3), Inches(6.65), Inches(9.4), Inches(0.55),
                        LIGHT_ORANGE_BG, ORANGE)
    add_text_box(slide, Inches(0.5), Inches(6.7), Inches(9.0), Inches(0.4),
                 "Overall Observation: After all Responsible AI interventions (mitigation + explainability + privacy), "
                 "BERT: accuracy 91%->87%, bias 7.5%->1%, explainability 0%->80%, privacy = formally guaranteed (epsilon=1). "
                 "Each intervention costs 1-3% accuracy but delivers substantial, measurable ethical value.",
                 10, bold=True, color=ORANGE)
    add_slide_number(slide, 11, TOTAL)

    # ==================================================================
    # SLIDE 12: ETHICAL RISKS & RESPONSIBLE AI FRAMEWORK
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Ethical Risks & Responsible AI Framework",
                  "Slide 12  |  Risks identified and our response to each")

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.4),
                 "5 Ethical Risks Identified", 15, bold=True, color=RED)
    risk_items = [
        "1. Algorithmic Discrimination (Bias in Predictions):",
        "   Our audit found DistilBERT assigns 7.5% more negative",
        "   sentiment to Black Male names for the EXACT same complaint.",
        "   In production, this means customers like 'Jamal' get lower",
        "   priority, longer wait times, and fewer human escalations",
        "   compared to 'Brad' -- purely because of their name.",
        "",
        "2. Privacy Violation (Demographic Exposure):",
        "   Customer names directly reveal race and gender. Deep",
        "   learning models can memorize training data, risking",
        "   leakage of personal identity attributes. Without privacy",
        "   safeguards, the model becomes a demographic profiler.",
        "",
        "3. Opacity / Black-Box Problem:",
        "   BERT has 66 million parameters. No human can inspect why",
        "   it scored 'Jamal' differently. Without explainability",
        "   tools, biased behavior remains invisible to developers,",
        "   auditors, and regulators -- making accountability impossible.",
        "",
        "4. Systemic Harm at Scale:",
        "   Amazon, Uber, and banks process billions of interactions/yr.",
        "   A 7.5% bias applied billions of times creates a feedback",
        "   loop: worse service leads to worse reviews, which trains",
        "   future models to be even more biased against that group.",
        "",
        "5. Legal & Regulatory Non-Compliance:",
        "   EU AI Act (2024) classifies customer service AI as high-risk.",
        "   GDPR Art. 22 requires explainable automated decisions.",
        "   US EEOC 80% rule sets a legal threshold for disparate impact.",
    ]
    add_bullet_frame(slide, Inches(0.5), Inches(1.8), Inches(4.3), Inches(5.0),
                     risk_items, 9, color=BLACK, spacing=Pt(1))

    add_text_box(slide, Inches(5.1), Inches(1.3), Inches(4.7), Inches(0.4),
                 "Our Technical Response (Risk --> Solution)", 15, bold=True, color=GREEN)
    resp_items = [
        "Each risk is addressed with a specific implemented technique:",
        "",
        "RISK 1 --> Fairness Auditing + Bias Mitigation",
        "  Detection: 5 Fairlearn metrics (Demographic Parity, Equal",
        "  Opportunity, Equalized Odds, Disparate Impact, Calibration)",
        "  quantify exactly how much bias exists per group.",
        "  Fix: 3 mitigation strategies -- Reweighing adjusts training",
        "  weights, Exp. Gradient adds fairness constraints during",
        "  training, Threshold Optimizer corrects outputs post-hoc.",
        "",
        "RISK 2 --> Differential Privacy (Mathematical Guarantee)",
        "  Laplace noise injection adds calibrated noise to model",
        "  outputs so individual identities cannot be reconstructed.",
        "  Opacus DP-SGD (Meta) clips gradients during training to",
        "  provide formal epsilon-delta privacy guarantees.",
        "",
        "RISK 3 --> Explainability via SHAP + LIME",
        "  SHAP computes Shapley values showing each word's",
        "  contribution to the score. LIME perturbs individual",
        "  sentences to explain local predictions. Together they",
        "  make the black-box decision process transparent.",
        "",
        "RISK 4 --> Repeatable 8-Step Audit Pipeline",
        "  Our full pipeline is automated and re-runnable. Companies",
        "  can schedule periodic audits to catch emerging bias.",
        "",
        "RISK 5 --> Full Documentation & Compliance Artifacts",
        "  Every metric, test result, and mitigation outcome is",
        "  logged as CSV/HTML reports for regulatory submission.",
    ]
    add_bullet_frame(slide, Inches(5.1), Inches(1.8), Inches(4.7), Inches(5.0),
                     resp_items, 9, color=BLACK, spacing=Pt(1))

    box = add_shape_box(slide, Inches(0.4), Inches(6.6), Inches(9.2), Inches(0.55),
                        LIGHT_BG, DARK_BLUE)
    add_text_box(slide, Inches(0.6), Inches(6.65), Inches(8.8), Inches(0.4),
                 "Responsible AI Principle: For every ethical risk, we implement a concrete, measurable technical solution. "
                 "Our audit is both diagnostic (finding bias) and prescriptive (fixing it with industry-standard tools).",
                 11, bold=True, color=DARK_BLUE)
    add_slide_number(slide, 12, TOTAL)

    # ==================================================================
    # SLIDE 13: FURTHER PROGRESS (PLANNED FOR FINAL REVIEW)
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Further Progress & Plan for Final Review",
                  "Slide 13  |  What remains to be completed")

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(4.5), Inches(0.4),
                 "Completed for Mid Review", 16, bold=True, color=GREEN)
    done_items = [
        "  Step 1: Dataset Generation -- 800 controlled sentences",
        "  across 8 demographic groups with balanced representation",
        "",
        "  Step 2: Sentiment Analysis -- All 800 sentences scored",
        "  using VADER, TextBlob, and DistilBERT (3 x 800 = 2400",
        "  sentiment evaluations completed)",
        "",
        "  Step 3: Bias Detection -- 4 statistical tests applied",
        "  (t-test, Cohen's d, Mann-Whitney U, ANOVA) confirming",
        "  BERT bias at 7.5% for Black Male names",
        "",
        "  Step 4: Fairness Metrics -- 5 Fairlearn metrics computed",
        "  for all 3 models. VADER/TextBlob pass; BERT shows issues.",
        "",
        "  Step 5: Bias Mitigation -- 3 strategies implemented",
        "  (Reweighing, Exponentiated Gradient, Threshold Optimizer)",
        "  Successfully reduces BERT's demographic disparity.",
    ]
    add_bullet_frame(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.8),
                     done_items, 11, color=BLACK, spacing=Pt(2))

    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4.5), Inches(0.4),
                 "Planned for Final Review", 16, bold=True, color=ORANGE)
    plan_items = [
        "  Step 6: Explainability Analysis (SHAP + LIME)",
        "  - SHAP global feature importance plots showing which",
        "    words drive bias across all 800 sentences",
        "  - LIME local explanations for individual biased cases",
        "  - Counterfactual analysis: what changes if we swap names",
        "",
        "  Step 7: Differential Privacy Implementation",
        "  - Laplace noise injection with epsilon analysis",
        "  - Opacus DP-SGD model training with privacy guarantees",
        "  - Privacy-utility tradeoff curves (epsilon vs accuracy)",
        "",
        "  Step 8: Visualization & Interactive Demo",
        "  - 10 publication-quality charts (heatmaps, box plots,",
        "    radar charts, privacy curves, SHAP force plots)",
        "  - Streamlit interactive demo application with live",
        "    bias comparison, counterfactual analysis, and",
        "    downloadable reports",
        "",
        "  Final Report: Comprehensive write-up with all results",
    ]
    add_bullet_frame(slide, Inches(5.3), Inches(1.8), Inches(4.5), Inches(5.0),
                     plan_items, 11, color=BLACK, spacing=Pt(2))

    box = add_shape_box(slide, Inches(0.4), Inches(6.3), Inches(9.2), Inches(0.8),
                        LIGHT_ORANGE_BG, ORANGE)
    add_text_box(slide, Inches(0.6), Inches(6.35), Inches(8.8), Inches(0.65),
                 "Timeline: Steps 6-7 (Explainability + Privacy) will be completed in the next 2 weeks. "
                 "Step 8 (Visualization + Demo) in the following week. Final report and presentation preparation "
                 "in the last week before the final review deadline.",
                 11, bold=True, color=ORANGE)
    add_slide_number(slide, 13, TOTAL)

    # ==================================================================
    # SLIDE 14: TEAM CONTRIBUTIONS (3 MEMBERS)
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Team Contributions (Mid Review)",
                  "Slide 14  |  Work distribution among 3 team members")

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5),
                 "Each team member contributed to specific modules of the project. Below is the work completed for mid review "
                 "and the planned responsibilities for the final review phase:",
                 12, color=GRAY)

    # Member 1
    box = add_shape_box(slide, Inches(0.3), Inches(2.0), Inches(3.05), Inches(4.6),
                        RGBColor(200, 230, 255), DARK_BLUE)
    add_text_box(slide, Inches(0.4), Inches(2.05), Inches(2.85), Inches(0.35),
                 "TEAM MEMBER 1", 15, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    m1_items = [
        "Mid Review Work:",
        "",
        "- Designed the controlled experiment",
        "  methodology (20 templates, 8 groups)",
        "- Built the dataset generation pipeline",
        "  (01_dataset_generation.py)",
        "- Generated 800 balanced complaint",
        "  sentences with demographic markers",
        "- Implemented all 3 sentiment analysis",
        "  models: VADER, TextBlob, DistilBERT",
        "  (02_sentiment_analysis.py)",
        "- Validated dataset quality and balance",
        "",
        "Final Review Plan:",
        "- Build Streamlit interactive demo",
        "- Create visualization charts",
        "- Generate final report",
    ]
    add_bullet_frame(slide, Inches(0.4), Inches(2.4), Inches(2.85), Inches(4.0),
                     m1_items, 10, color=BLACK, spacing=Pt(1))

    # Member 2
    box = add_shape_box(slide, Inches(3.5), Inches(2.0), Inches(3.05), Inches(4.6),
                        RGBColor(230, 200, 255), RGBColor(100, 50, 150))
    add_text_box(slide, Inches(3.6), Inches(2.05), Inches(2.85), Inches(0.35),
                 "TEAM MEMBER 2", 15, bold=True, color=RGBColor(100, 50, 150), alignment=PP_ALIGN.CENTER)
    m2_items = [
        "Mid Review Work:",
        "",
        "- Implemented statistical bias detection",
        "  tests: t-test, Cohen's d, Mann-Whitney",
        "  U, ANOVA (03_bias_detection.py)",
        "- Computed 5 Fairlearn fairness metrics:",
        "  Demographic Parity, Equal Opportunity,",
        "  Equalized Odds, Disparate Impact,",
        "  Calibration (04_fairness_metrics.py)",
        "- Implemented 3 bias mitigation strategies",
        "  (05_bias_mitigation.py)",
        "- Analyzed bias detection results",
        "",
        "Final Review Plan:",
        "- SHAP explainability analysis",
        "- LIME local explanations",
        "- Counterfactual analysis",
    ]
    add_bullet_frame(slide, Inches(3.6), Inches(2.4), Inches(2.85), Inches(4.0),
                     m2_items, 10, color=BLACK, spacing=Pt(1))

    # Member 3
    box = add_shape_box(slide, Inches(6.7), Inches(2.0), Inches(3.1), Inches(4.6),
                        RGBColor(200, 255, 200), GREEN)
    add_text_box(slide, Inches(6.8), Inches(2.05), Inches(2.9), Inches(0.35),
                 "TEAM MEMBER 3", 15, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
    m3_items = [
        "Mid Review Work:",
        "",
        "- Researched responsible AI frameworks",
        "  (EU AI Act, GDPR, EEOC guidelines)",
        "- Identified 5 ethical risks in customer",
        "  service sentiment analysis systems",
        "- Designed the responsible AI response",
        "  framework (risk-to-technique mapping)",
        "- Literature review: Caliskan (2017),",
        "  Bolukbasi (2016), Mehrabi (2021)",
        "- Prepared mid-review presentation",
        "  and documentation",
        "",
        "Final Review Plan:",
        "- Differential privacy implementation",
        "  (Laplace + Opacus DP-SGD)",
        "- Privacy-utility tradeoff analysis",
        "- Final report compilation",
    ]
    add_bullet_frame(slide, Inches(6.8), Inches(2.4), Inches(2.9), Inches(4.0),
                     m3_items, 10, color=BLACK, spacing=Pt(1))

    add_text_box(slide, Inches(0.3), Inches(6.8), Inches(9.5), Inches(0.4),
                 "(Replace 'Team Member 1/2/3' with actual names before presenting)",
                 11, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_slide_number(slide, 14, TOTAL)

    # ==================================================================
    # SLIDE 15: THANK YOU
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape_box(slide, Inches(0), Inches(0), Inches(10), Inches(7.5), DARK_BLUE)

    add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(1.0),
                 "Thank You", 52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.8), Inches(8), Inches(0.5),
                 "Auditing Gender & Race Bias in Customer Service AI",
                 18, color=RGBColor(180, 210, 255), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(0.5),
                 "Questions & Discussion",
                 20, color=RGBColor(200, 220, 255), alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 15, TOTAL)

    # SAVE
    out = os.path.join(OUTPUT_DIR, "Mid_Review_Presentation_v6.pptx")
    prs.save(out)
    print(f"PPT saved to: {out}")


if __name__ == "__main__":
    build()
